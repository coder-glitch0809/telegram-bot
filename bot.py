import asyncio
import base64
import html
import logging
import os
import re
import shutil
import smtplib
import sqlite3
import secrets
import string
import tempfile
import time
import uuid
from datetime import datetime, timedelta
from email.message import EmailMessage
from pathlib import Path
from typing import Any, cast
from urllib.parse import urlparse

import httpx
from dotenv import load_dotenv
from fastapi import FastAPI, HTTPException, Request
from openai import AsyncOpenAI, OpenAI
from telegram import BotCommand, InlineKeyboardButton, InlineKeyboardMarkup, MessageEntity, Update
from telegram.constants import ChatAction, ChatType, ParseMode
from telegram.error import BadRequest, NetworkError, TimedOut
from telegram.ext import (
    Application,
    CallbackQueryHandler,
    CommandHandler,
    ContextTypes,
    MessageHandler,
    filters,
)

Agent = Runner = None
set_default_openai_api = set_default_openai_client = set_tracing_disabled = None
try:
    from agents import Agent, Runner, set_default_openai_api, set_default_openai_client, set_tracing_disabled
except ImportError:
    pass


load_dotenv()

TELEGRAM_BOT_TOKEN = os.getenv("TELEGRAM_BOT_TOKEN", os.getenv("TELEGRAM_BOT_token", "")).strip()
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "").strip()
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
OPENAI_IMAGE_API_KEY = os.getenv("OPENAI_IMAGE_API_KEY", OPENAI_API_KEY).strip()

AI_PROVIDER_ENV = os.getenv("AI_PROVIDER", "").strip().lower()
SUPPORTED_AI_PROVIDERS = {"groq", "openai"}
if AI_PROVIDER_ENV in SUPPORTED_AI_PROVIDERS:
    AI_PROVIDER = AI_PROVIDER_ENV
elif GROQ_API_KEY:
    AI_PROVIDER = "groq"
else:
    AI_PROVIDER = "openai"

AI_API_KEY = {
    "groq": GROQ_API_KEY,
    "openai": OPENAI_API_KEY,
}.get(AI_PROVIDER, OPENAI_API_KEY)

DEFAULT_AI_BASE_URL = {
    "groq": "https://api.groq.com/openai/v1",
    "openai": "",
}.get(AI_PROVIDER, "")
AI_BASE_URL = os.getenv("AI_BASE_URL", DEFAULT_AI_BASE_URL).strip()
if AI_PROVIDER == "groq" and AI_BASE_URL.rstrip("/") in {"https://api.groq.com/v1", "https://api.groq.com/v1/models"}:
    AI_BASE_URL = "https://api.groq.com/openai/v1"

DEFAULT_TEXT_MODEL = {
    "groq": "llama-3.3-70b-versatile",
    "openai": "gpt-4o-mini",
}.get(AI_PROVIDER, "gpt-4o-mini")
OPENAI_TEXT_MODEL = os.getenv("OPENAI_TEXT_MODEL", DEFAULT_TEXT_MODEL).strip()
DEFAULT_TRANSCRIBE_MODEL = {
    "groq": "whisper-large-v3",
    "openai": "gpt-4o-mini-transcribe",
}.get(AI_PROVIDER, "gpt-4o-mini-transcribe")
OPENAI_TRANSCRIBE_MODEL = os.getenv("OPENAI_TRANSCRIBE_MODEL", DEFAULT_TRANSCRIBE_MODEL).strip()

OWNER_TELEGRAM_ID = int(os.getenv("OWNER_TELEGRAM_ID", os.getenv("BOT_OWNER_ID", "0")) or 0)
OWNER_EMAIL = os.getenv("OWNER_EMAIL", "").strip()
SMTP_HOST = os.getenv("SMTP_HOST", "").strip()
SMTP_PORT = int(os.getenv("SMTP_PORT", "587") or 587)
SMTP_USERNAME = os.getenv("SMTP_USERNAME", "").strip()
SMTP_PASSWORD = os.getenv("SMTP_PASSWORD", "").strip()
SMTP_FROM_EMAIL = os.getenv("SMTP_FROM_EMAIL", SMTP_USERNAME).strip()
REPORT_WEEKLY_DAY = int(os.getenv("REPORT_WEEKLY_DAY", "0") or 0)

IMAGE_GENERATION_ENABLED = os.getenv("IMAGE_GENERATION_ENABLED", "true").strip().lower() in {"1", "true", "yes", "ha"}
IMAGE_MODEL = os.getenv("IMAGE_MODEL", "gpt-image-1").strip()
IMAGE_SIZE = os.getenv("IMAGE_SIZE", "1024x1024").strip()
START_EMOJI = "\U0001F44B"
START_CUSTOM_EMOJI_ID = "5472055112702629499"
AUDIO_BUTTON_CUSTOM_EMOJI_ID = "5260325873688518261"
VIDEO_BUTTON_CUSTOM_EMOJI_ID = "5375309569905938163"
MEDIA_DONE_EMOJI = "\u2714\ufe0f"
MEDIA_DONE_CUSTOM_EMOJI_ID = "5321210956414459578"
WARNING_EMOJI = "\u26a0\ufe0f"
WARNING_CUSTOM_EMOJI_ID = "5215305227627931680"
IDEA_EMOJI = "\U0001F4A1"
IDEA_CUSTOM_EMOJI_ID = "5355014749920709843"
BOT_EMOJI = "\U0001F916"
BOT_CUSTOM_EMOJI_ID = "5372981976804366741"
PHOTO_EMOJI = "\U0001F4F7"
PHOTO_CUSTOM_EMOJI_ID = "5429187662596561219"
CHAT_EMOJI = "\U0001F4AC"
CHAT_CUSTOM_EMOJI_ID = "5443038326535759644"
SEARCH_EMOJI = "\U0001F50E"
SEARCH_CUSTOM_EMOJI_ID = "5188311512791393083"
WAIT_EMOJI = "\u267e\ufe0f"
WAIT_CUSTOM_EMOJI_ID = "5389019921558563669"

MEDIA_DOWNLOAD_ENABLED = os.getenv("MEDIA_DOWNLOAD_ENABLED", os.getenv("YOUTUBE_DOWNLOAD_ENABLED", "true")).strip().lower() in {
    "1",
    "true",
    "yes",
    "ha",
}
MEDIA_MAX_MB = int(os.getenv("MEDIA_MAX_MB", os.getenv("YOUTUBE_MAX_MB", "45")) or 45)
KINOTOP_OWNER_ID = 896778319
LOVE_OWNER_ID = 896778319
LOVE_CODE_ALPHABET = string.ascii_uppercase + string.digits


def normalize_domain(value: str) -> str:
    value = value.strip().lower().rstrip("/")
    if not value:
        return ""
    parsed = urlparse(value if "://" in value else f"https://{value}")
    return (parsed.hostname or "").lower().removeprefix("www.")


KINOTOP_ALLOWED_DOMAINS = {
    domain
    for domain in (normalize_domain(item) for item in os.getenv("KINOTOP_ALLOWED_DOMAINS", "").split(","))
    if domain
}
DEFAULT_ANALYTICS_DB_FILE = (
    str(Path(tempfile.gettempdir()) / "bot_analytics.sqlite3")
    if os.getenv("VERCEL") or os.getenv("VERCEL_URL")
    else "bot_analytics.sqlite3"
)
ANALYTICS_DB_FILE = os.getenv("ANALYTICS_DB_FILE", DEFAULT_ANALYTICS_DB_FILE).strip()
if os.getenv("VERCEL") or os.getenv("VERCEL_URL"):
    ANALYTICS_DB_FILE = str(Path(tempfile.gettempdir()) / "bot_analytics.sqlite3")
DEFAULT_TEMP_OUTPUT_DIR = (
    str(Path(tempfile.gettempdir()))
    if os.getenv("VERCEL") or os.getenv("VERCEL_URL")
    else ".bot-temp"
)
TEMP_OUTPUT_DIR = os.getenv("BOT_TEMP_DIR", DEFAULT_TEMP_OUTPUT_DIR).strip()

URL_RE = re.compile(r"https?://\S+", re.IGNORECASE)
ANSI_RE = re.compile(r"\x1b\[[0-9;]*m")
HTML_TAG_RE = re.compile(r"<[^>]+>")
MEDIA_URL_RE = re.compile(
    r"^https?://(www\.)?(youtube\.com|youtu\.be|music\.youtube\.com|instagram\.com|instagr\.am)/\S+$",
    re.IGNORECASE,
)
EXCEL_LIST_RE = re.compile(
    r"\b(excel|xlsx|csv|tsv|jadval|table|spreadsheet|spisok|spiska|ro'yxat|royxat|list)\b",
    re.IGNORECASE,
)
IMAGE_WORDS = {"rasm", "surat", "image", "picture", "нарисуй", "изображение", "сгенерируй"}
PRESENTATION_WORDS = {
    "present",
    "presentation",
    "prezentatsiya",
    "презентация",
    "slayd",
    "slide",
    "ppt",
    "pptx",
    "dars ishlanma",
    "konspekt",
    "referat",
}
ADULT_WORDS = {
    "18+",
    "porn",
    "porno",
    "xxx",
    "nsfw",
    "nude",
    "naked",
    "yalangoch",
    "yalang'och",
    "pornografiya",
    "порно",
    "нюд",
    "голый",
    "голая",
    "эротика",
}

logging.basicConfig(format="%(asctime)s - %(name)s - %(levelname)s - %(message)s", level=logging.INFO)
logger = logging.getLogger(__name__)
logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("httpcore").setLevel(logging.WARNING)

if AI_BASE_URL:
    openai_client = OpenAI(api_key=AI_API_KEY or "missing", base_url=AI_BASE_URL)
    agents_client = AsyncOpenAI(api_key=AI_API_KEY or "missing", base_url=AI_BASE_URL)
else:
    openai_client = OpenAI(api_key=AI_API_KEY or "missing")
    agents_client = AsyncOpenAI(api_key=AI_API_KEY or "missing")
if set_default_openai_client and set_default_openai_api and set_tracing_disabled:
    set_default_openai_client(agents_client, use_for_tracing=False)
    set_default_openai_api("chat_completions")
    set_tracing_disabled(True)
image_client = OpenAI(api_key=OPENAI_IMAGE_API_KEY or "missing")
analytics: "AnalyticsStore | None" = None
bot_app: "Application | None" = None
bot_lock = asyncio.Lock()
webhook_lock = asyncio.Lock()
webhook_ready = False
bot_init_error = ""


class AnalyticsStore:
    def __init__(self, db_file: str) -> None:
        self.db_file = db_file
        self._init_db()

    def _connect(self) -> sqlite3.Connection:
        db_path = Path(self.db_file)
        if db_path.parent and str(db_path.parent) not in {"", "."}:
            db_path.parent.mkdir(parents=True, exist_ok=True)
        connection = sqlite3.connect(self.db_file)
        connection.row_factory = sqlite3.Row
        return connection

    def _ensure_column(self, connection: sqlite3.Connection, table: str, column: str, definition: str) -> None:
        columns = {row["name"] for row in connection.execute(f"PRAGMA table_info({table})").fetchall()}
        if column not in columns:
            connection.execute(f"ALTER TABLE {table} ADD COLUMN {column} {definition}")

    def _init_db(self) -> None:
        with self._connect() as connection:
            connection.execute(
                """
                CREATE TABLE IF NOT EXISTS users (
                    user_id INTEGER PRIMARY KEY,
                    username TEXT,
                    full_name TEXT,
                    first_seen TEXT NOT NULL,
                    last_seen TEXT NOT NULL,
                    message_count INTEGER NOT NULL DEFAULT 0,
                    ai_count INTEGER NOT NULL DEFAULT 0,
                    voice_count INTEGER NOT NULL DEFAULT 0,
                    image_count INTEGER NOT NULL DEFAULT 0,
                    media_count INTEGER NOT NULL DEFAULT 0
                )
                """
            )
            self._ensure_column(connection, "users", "image_count", "INTEGER NOT NULL DEFAULT 0")
            self._ensure_column(connection, "users", "media_count", "INTEGER NOT NULL DEFAULT 0")
            connection.execute(
                """
                CREATE TABLE IF NOT EXISTS interactions (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    created_at TEXT NOT NULL,
                    user_id INTEGER NOT NULL,
                    username TEXT,
                    action TEXT NOT NULL,
                    status TEXT NOT NULL,
                    text_preview TEXT,
                    response_chars INTEGER NOT NULL DEFAULT 0,
                    error TEXT
                )
                """
            )
            connection.execute(
                """
                CREATE TABLE IF NOT EXISTS sent_reports (
                    report_key TEXT PRIMARY KEY,
                    sent_at TEXT NOT NULL
                )
                """
            )
            connection.execute(
                """
                CREATE TABLE IF NOT EXISTS love_sessions (
                    code TEXT PRIMARY KEY,
                    owner_id INTEGER NOT NULL,
                    partner_id INTEGER,
                    partner_name TEXT,
                    created_at TEXT NOT NULL,
                    joined_at TEXT,
                    closed_at TEXT
                )
                """
            )

    def record_interaction(
        self,
        user_id: int,
        username: str,
        full_name: str,
        action: str,
        status: str = "ok",
        text_preview: str = "",
        response_chars: int = 0,
        error: str = "",
    ) -> None:
        now = datetime.now().isoformat(timespec="seconds")
        preview = text_preview.replace("\n", " ")[:300]
        with self._connect() as connection:
            connection.execute(
                """
                INSERT INTO users (user_id, username, full_name, first_seen, last_seen)
                VALUES (?, ?, ?, ?, ?)
                ON CONFLICT(user_id) DO UPDATE SET
                    username = excluded.username,
                    full_name = excluded.full_name,
                    last_seen = excluded.last_seen
                """,
                (user_id, username, full_name, now, now),
            )
            connection.execute(
                """
                UPDATE users
                SET
                    message_count = message_count + 1,
                    ai_count = ai_count + ?,
                    voice_count = voice_count + ?,
                    image_count = image_count + ?,
                    media_count = media_count + ?
                WHERE user_id = ?
                """,
                (
                    1 if action in {"ai", "text"} and status == "ok" else 0,
                    1 if action == "voice" and status == "ok" else 0,
                    1 if action == "image" and status == "ok" else 0,
                    1 if action == "media" and status == "ok" else 0,
                    user_id,
                ),
            )
            connection.execute(
                """
                INSERT INTO interactions (created_at, user_id, username, action, status, text_preview, response_chars, error)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                """,
                (now, user_id, username, action, status, preview, response_chars, error[:300]),
            )

    def user_image_count(self, user_id: int) -> int:
        with self._connect() as connection:
            row = connection.execute("SELECT image_count FROM users WHERE user_id = ?", (user_id,)).fetchone()
        return int(row["image_count"] or 0) if row else 0

    def user_exists(self, user_id: int) -> bool:
        with self._connect() as connection:
            row = connection.execute("SELECT 1 FROM users WHERE user_id = ?", (user_id,)).fetchone()
        return row is not None

    def has_sent_report(self, report_key: str) -> bool:
        with self._connect() as connection:
            row = connection.execute("SELECT 1 FROM sent_reports WHERE report_key = ?", (report_key,)).fetchone()
        return row is not None

    def mark_report_sent(self, report_key: str) -> None:
        with self._connect() as connection:
            connection.execute(
                "INSERT OR REPLACE INTO sent_reports (report_key, sent_at) VALUES (?, ?)",
                (report_key, datetime.now().isoformat(timespec="seconds")),
            )

    def create_love_invite(self, owner_id: int, code: str) -> None:
        now = datetime.now().isoformat(timespec="seconds")
        with self._connect() as connection:
            connection.execute(
                "UPDATE love_sessions SET closed_at = ? WHERE owner_id = ? AND closed_at IS NULL",
                (now, owner_id),
            )
            connection.execute(
                """
                INSERT INTO love_sessions (code, owner_id, created_at)
                VALUES (?, ?, ?)
                """,
                (code, owner_id, now),
            )

    def join_love_session(self, code: str, partner_id: int, partner_name: str) -> sqlite3.Row | None:
        now = datetime.now().isoformat(timespec="seconds")
        with self._connect() as connection:
            row = connection.execute(
                """
                SELECT * FROM love_sessions
                WHERE code = ? AND closed_at IS NULL
                """,
                (code,),
            ).fetchone()
            if row is None or int(row["owner_id"]) == partner_id:
                return None
            if row["partner_id"] and int(row["partner_id"]) != partner_id:
                return row
            connection.execute(
                """
                UPDATE love_sessions
                SET partner_id = ?, partner_name = ?, joined_at = COALESCE(joined_at, ?)
                WHERE code = ?
                """,
                (partner_id, partner_name[:120], now, code),
            )
            return connection.execute("SELECT * FROM love_sessions WHERE code = ?", (code,)).fetchone()

    def get_love_session_by_user(self, user_id: int) -> sqlite3.Row | None:
        with self._connect() as connection:
            return connection.execute(
                """
                SELECT * FROM love_sessions
                WHERE closed_at IS NULL
                AND (owner_id = ? OR partner_id = ?)
                ORDER BY created_at DESC
                LIMIT 1
                """,
                (user_id, user_id),
            ).fetchone()

    def close_love_session(self, user_id: int) -> sqlite3.Row | None:
        now = datetime.now().isoformat(timespec="seconds")
        with self._connect() as connection:
            row = self.get_love_session_by_user(user_id)
            if row is None:
                return None
            connection.execute("UPDATE love_sessions SET closed_at = ? WHERE code = ?", (now, row["code"]))
            return row

    def subscriber_count(self) -> dict[str, int]:
        now = datetime.now()
        day_ago = (now - timedelta(days=1)).isoformat(timespec="seconds")
        week_ago = (now - timedelta(days=7)).isoformat(timespec="seconds")
        with self._connect() as connection:
            row = connection.execute(
                """
                SELECT
                    COUNT(*) AS total,
                    SUM(CASE WHEN first_seen >= ? THEN 1 ELSE 0 END) AS new_today,
                    SUM(CASE WHEN first_seen >= ? THEN 1 ELSE 0 END) AS new_week,
                    SUM(CASE WHEN last_seen >= ? THEN 1 ELSE 0 END) AS active_today,
                    SUM(CASE WHEN last_seen >= ? THEN 1 ELSE 0 END) AS active_week
                FROM users
                """,
                (day_ago, week_ago, day_ago, week_ago),
            ).fetchone()
        return {
            "total": int(row["total"] or 0),
            "new_today": int(row["new_today"] or 0),
            "new_week": int(row["new_week"] or 0),
            "active_today": int(row["active_today"] or 0),
            "active_week": int(row["active_week"] or 0),
        }

    def top_queries(self, days: int = 7, limit: int = 10) -> list[str]:
        start_at = (datetime.now() - timedelta(days=days)).isoformat(timespec="seconds")
        with self._connect() as connection:
            rows = connection.execute(
                """
                SELECT text_preview, COUNT(*) AS count
                FROM interactions
                WHERE created_at >= ?
                AND text_preview != ''
                AND action IN ('ai', 'text', 'voice', 'image', 'media')
                AND status = 'ok'
                GROUP BY LOWER(text_preview)
                ORDER BY count DESC, MAX(created_at) DESC
                LIMIT ?
                """,
                (start_at, limit),
            ).fetchall()
        return [f"{row['text_preview']} ({row['count']}x)" for row in rows]

    def build_report(self, title: str, start_at: datetime | None = None, end_at: datetime | None = None) -> str:
        conditions = []
        params_list: list[Any] = []
        if start_at:
            conditions.append("created_at >= ?")
            params_list.append(start_at.isoformat(timespec="seconds"))
        if end_at:
            conditions.append("created_at < ?")
            params_list.append(end_at.isoformat(timespec="seconds"))
        where = f"WHERE {' AND '.join(conditions)}" if conditions else ""
        params = tuple(params_list)

        with self._connect() as connection:
            total_users = connection.execute("SELECT COUNT(*) FROM users").fetchone()[0]
            active_users = connection.execute(f"SELECT COUNT(DISTINCT user_id) FROM interactions {where}", params).fetchone()[0]
            action_rows = connection.execute(
                f"SELECT action, COUNT(*) AS count FROM interactions {where} GROUP BY action ORDER BY count DESC",
                params,
            ).fetchall()
            top_rows = connection.execute(
                f"""
                SELECT text_preview, COUNT(*) AS count
                FROM interactions
                {where}
                GROUP BY LOWER(text_preview)
                HAVING text_preview != ''
                ORDER BY count DESC, MAX(created_at) DESC
                LIMIT 12
                """,
                params,
            ).fetchall()

        lines = [
            title,
            "",
            f"Jami obunachi/foydalanuvchi: {total_users}",
            f"Shu davrda faol foydalanuvchi: {active_users}",
            "",
            "Ishlatish turi:",
        ]
        lines.extend(f"- {row['action']}: {row['count']}" for row in action_rows) if action_rows else lines.append("- Hali faollik yo'q")
        lines.extend(["", "Kim ko'p nima qidiryapti / so'rayapti:"])
        lines.extend(f"- {row['text_preview']} ({row['count']}x)" for row in top_rows) if top_rows else lines.append("- Hali yetarli so'rov yo'q")
        return "\n".join(lines)


def require_config() -> None:
    missing = []
    if not TELEGRAM_BOT_TOKEN:
        missing.append("TELEGRAM_BOT_TOKEN")
    if not AI_API_KEY:
        missing.append("GROQ_API_KEY yoki OPENAI_API_KEY")
    if missing:
        raise RuntimeError(f"Sozlanmagan yoki topilmadi: {', '.join(missing)}. .env faylini ko'ring.")


def get_analytics() -> AnalyticsStore:
    global analytics
    if analytics is None:
        analytics = AnalyticsStore(ANALYTICS_DB_FILE)
    return analytics


def is_owner(user_id: int) -> bool:
    return OWNER_TELEGRAM_ID != 0 and user_id == OWNER_TELEGRAM_ID


def make_love_code(length: int = 8) -> str:
    return "".join(secrets.choice(LOVE_CODE_ALPHABET) for _ in range(length))


def get_url_domain(url: str) -> str:
    return normalize_domain(url)


def is_kinotop_allowed_url(url: str) -> bool:
    domain = get_url_domain(url)
    if not domain or not KINOTOP_ALLOWED_DOMAINS:
        return False
    return any(domain == allowed or domain.endswith(f".{allowed}") for allowed in KINOTOP_ALLOWED_DOMAINS)


def email_reports_enabled() -> bool:
    return all([OWNER_EMAIL, SMTP_HOST, SMTP_USERNAME, SMTP_PASSWORD, SMTP_FROM_EMAIL])


def send_email(subject: str, body: str) -> None:
    if not email_reports_enabled():
        return
    message = EmailMessage()
    message["Subject"] = subject
    message["From"] = SMTP_FROM_EMAIL
    message["To"] = OWNER_EMAIL
    message.set_content(body)
    with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=30) as smtp:
        smtp.starttls()
        smtp.login(SMTP_USERNAME, SMTP_PASSWORD)
        smtp.send_message(message)


async def send_owner_report(application: Application, subject: str, body: str) -> None:
    if OWNER_TELEGRAM_ID:
        await application.bot.send_message(chat_id=OWNER_TELEGRAM_ID, text=body[:4096])
    await asyncio.to_thread(send_email, subject, body)


async def maybe_send_reports(application: Application | None = None) -> None:
    now = datetime.now()
    if now.weekday() != REPORT_WEEKLY_DAY:
        return
    store = get_analytics()
    report_key = f"weekly-{now.strftime('%Y-%W')}"
    if store.has_sent_report(report_key):
        return
    body = store.build_report("Haftalik Telegram AI bot hisoboti", start_at=now - timedelta(days=7))
    if application:
        await send_owner_report(application, "Haftalik Telegram AI bot hisoboti", body)
    else:
        await asyncio.to_thread(send_email, "Haftalik Telegram AI bot hisoboti", body)
    store.mark_report_sent(report_key)


async def report_scheduler(application: Application) -> None:
    while True:
        try:
            await maybe_send_reports(application)
        except Exception:
            logger.exception("Haftalik hisobot yuborishda xato")
        await asyncio.sleep(60 * 60)


async def start_background_tasks(application: Application) -> None:
    await setup_bot_commands(application)
    application.bot_data["report_scheduler_task"] = asyncio.create_task(report_scheduler(application))


async def stop_background_tasks(application: Application) -> None:
    task = application.bot_data.get("report_scheduler_task")
    if task:
        task.cancel()
        try:
            await task
        except asyncio.CancelledError:
            pass


async def setup_bot_commands(application: Application) -> None:
    await application.bot.set_my_commands(
        [
            BotCommand("start", "Botni boshlash"),
            BotCommand("ai", "Guruhda yoki privatda AI savol"),
            BotCommand("image", "AI rasm generatsiyasi"),
            BotCommand("present", "Prezentatsiya yoki o'quv fayli"),
            BotCommand("media", "Instagram/YouTube audio yoki video"),
            BotCommand("radar", "Foydalanuvchilar va trendlar"),
            BotCommand("report", "Admin haftalik hisobot"),
            BotCommand("help", "Yordam"),
        ]
    )


def require_effective_user(update: Update) -> Any:
    user = update.effective_user
    if user is None:
        raise RuntimeError("Telegram update ichida foydalanuvchi topilmadi.")
    return user


def require_message(update: Update) -> Any:
    message = update.message or update.effective_message
    if message is None:
        raise RuntimeError("Telegram update ichida xabar topilmadi.")
    return message


def make_temp_dir(prefix: str) -> Path:
    base_dir = Path(TEMP_OUTPUT_DIR)
    base_dir.mkdir(parents=True, exist_ok=True)
    for _ in range(10):
        path = base_dir / f"{prefix}{uuid.uuid4().hex[:10]}"
        try:
            path.mkdir(parents=True, exist_ok=False)
            return path
        except FileExistsError:
            continue
    raise RuntimeError("Vaqtinchalik papka yaratib bo'lmadi.")


def context_args(context: ContextTypes.DEFAULT_TYPE) -> list[str]:
    return list(context.args or [])


def telegram_user_details(update: Update) -> tuple[int, str, str]:
    user = require_effective_user(update)
    username = user.username or ""
    full_name = user.full_name or username or "user"
    return user.id, username, full_name


def record_usage(update: Update, action: str, status: str = "ok", text_preview: str = "", response_chars: int = 0, error: str = "") -> None:
    user_id, username, full_name = telegram_user_details(update)
    try:
        get_analytics().record_interaction(user_id, username, full_name, action, status, text_preview, response_chars, error)
    except Exception:
        logger.exception("Analytics yozishda xato")


def is_group_chat(update: Update) -> bool:
    chat = update.effective_chat
    return chat is not None and chat.type in {ChatType.GROUP, ChatType.SUPERGROUP}


def contains_adult_content(text: str) -> bool:
    lowered = text.lower()
    return any(word in lowered for word in ADULT_WORDS)


def looks_like_image_request(text: str) -> bool:
    lowered = text.lower()
    return any(word in lowered for word in IMAGE_WORDS)


def looks_like_presentation_request(text: str) -> bool:
    lowered = text.lower()
    return any(word in lowered for word in PRESENTATION_WORDS)


def looks_like_excel_list_request(text: str) -> bool:
    return bool(EXCEL_LIST_RE.search(text))


def is_transient_ai_error(exc: Exception) -> bool:
    message = str(exc).lower()
    return any(
        token in message
        for token in ["503", "unavailable", "high demand", "temporarily unavailable", "rate limit", "timeout", "connection reset"]
    )


def ai_call_with_retries(callable_obj, retries: int = 3, initial_delay: float = 1.0):
    delay = initial_delay
    for attempt in range(1, retries + 1):
        try:
            return callable_obj()
        except Exception as exc:
            if attempt == retries or not is_transient_ai_error(exc):
                raise
            logger.warning("AI request transient error, retrying %d/%d after %.1fs: %s", attempt, retries, delay, exc)
            time.sleep(delay)
            delay = min(delay * 2, 5.0)


def friendly_error(exc: Exception) -> str:
    message = ANSI_RE.sub("", str(exc))
    lowered = message.lower()
    if "invalid_api_key" in lowered or "invalid api key" in lowered or "401" in message:
        expected_key = {
            "groq": "GROQ_API_KEY",
            "openai": "OPENAI_API_KEY",
        }.get(AI_PROVIDER, "AI API key")
        return (
            f"AI API key noto'g'ri yoki providerga mos emas. Hozirgi provider: {AI_PROVIDER}. "
            f".env ichida {expected_key} ni yangilang, AI_BASE_URL ni providerga mos qiling va botni qayta ishga tushiring."
        )
    if "requested format is not available" in lowered:
        return "Bu media uchun so'ralgan format topilmadi. Boshqa formatni tanlab ko'ring."
    if "video unavailable" in lowered:
        return "Bu video hozir mavjud emas yoki yopiq."
    if "timed out" in lowered or "timeout" in lowered:
        return "Media tayyorlash yoki yuborish vaqti tugadi. Qisqaroq link yuboring yoki qayta urinib ko'ring."
    if "ffmpeg" in lowered:
        return "Serverda ffmpeg yo'q. Video uchun bitta tayyor format tanlab ko'ring yoki ffmpeg o'rnating."
    if "permission" in lowered or "403" in message:
        return "AI xizmatiga kirish ruxsati yo'q. API key yoki billingni tekshiring."
    if "503" in message or "unavailable" in lowered or "high demand" in lowered:
        return "AI modeli hozir band. Birozdan keyin qayta urinib ko'ring."
    if "connect" in lowered or "connection" in lowered:
        return "Internet yoki API ulanishida muammo bor. Birozdan keyin qayta urinib ko'ring."
    return message[:700]


def build_chat_instructions(text: str) -> str:
    excel_mode = looks_like_excel_list_request(text)
    emoji_guide = (
        f"Bot/yordamchi mavzusi: {custom_emoji_html(BOT_EMOJI, BOT_CUSTOM_EMOJI_ID)}; "
        f"ogohlantirish/xavf: {custom_emoji_html(WARNING_EMOJI, WARNING_CUSTOM_EMOJI_ID)}; "
        f"fikr/maslahat: {custom_emoji_html(IDEA_EMOJI, IDEA_CUSTOM_EMOJI_ID)}; "
        f"rasm/kamera: {custom_emoji_html(PHOTO_EMOJI, PHOTO_CUSTOM_EMOJI_ID)}; "
        f"savol-javob/chat: {custom_emoji_html(CHAT_EMOJI, CHAT_CUSTOM_EMOJI_ID)}; "
        f"qidirish/tahlil: {custom_emoji_html(SEARCH_EMOJI, SEARCH_CUSTOM_EMOJI_ID)}; "
        f"kutish/jarayon: {custom_emoji_html(WAIT_EMOJI, WAIT_CUSTOM_EMOJI_ID)}."
    )
    return (
        "Siz Telegram ichidagi zamonaviy AI yordamchisiz. "
        "Asosan o'zbekcha, ruscha va inglizcha muloqot qiling; foydalanuvchi qaysi tilda yozsa, shu tilda javob bering. "
        "Javoblar aniq, foydali, tabiiy va qisqa bo'lsin. "
        "Siz maktab o'qituvchilari, repetitorlar, talabalar va o'quvchilarga dars reja, konspekt, test, izoh, "
        "misol, uy vazifasi uchun yo'nalish, prezentatsiya matni, jadval va o'quv materiallarini tayyorlashda yordam berasiz. "
        "O'quvchiga tayyor ko'chirma emas, tushunarli bosqichma-bosqich yechim va o'rganishga yordam beradigan izoh bering. "
        "Ustozlar uchun dars maqsadi, metod, vaqt taqsimoti, baholash mezoni va topshiriqlarni tartibli tuzing. "
        "Javobni Telegram HTML formatida yozing: faqat <b>, <i>, <code>, <pre> va <tg-emoji> taglaridan foydalaning. "
        "Markdown belgilarini ishlatmang. Mavzuga mos 1-4 ta animated emoji qo'shing, lekin ortiqcha bezamang. "
        f"Faqat shu custom emoji taglaridan foydalaning: {emoji_guide} "
        "18+ pornografik, jinsiy ekspluatatsiya yoki noqonuniy materiallarni yaratmang, topmang va tarqatmang. "
        "Bunday so'rovda qisqa rad etib, xavfsiz alternativ taklif qiling. "
        "HTML taglarni doim to'g'ri yoping va javobni Telegramda o'qishga qulay qiling."
        + (
            " Agar foydalanuvchi Excel, CSV, jadval, ro'yxat, spisok yoki spreadsheet uchun format so'rasa, "
            "javobni faqat <pre> ichida tab bilan ajratilgan TSV jadval ko'rinishida yozing. "
            "Birinchi qatorda ustun nomlari bo'lsin. Emoji, markdown, ortiqcha izoh va HTML tag ishlatmang; faqat bitta <pre>...</pre> blok qaytaring."
            if excel_mode
            else ""
        )
    )


async def ask_ai(text: str) -> str:
    if not AI_API_KEY:
        raise RuntimeError(f"{AI_PROVIDER.upper()} uchun API key topilmadi. .env ichidagi AI_PROVIDER va mos API keyni tekshiring.")

    instructions = build_chat_instructions(text)
    if Agent and Runner:
        chat_agent = Agent(
            name="Telegram Chat Agent",
            instructions=instructions,
            model=OPENAI_TEXT_MODEL,
        )
        result = await Runner.run(chat_agent, text)
        return str(result.final_output or "").strip()

    response = cast(Any, await asyncio.to_thread(
        lambda: ai_call_with_retries(
            lambda: openai_client.chat.completions.create(
                model=OPENAI_TEXT_MODEL,
                messages=[
                    {"role": "system", "content": instructions},
                    {"role": "user", "content": text},
                ],
                temperature=0.45,
            )
        )
    ))
    return (response.choices[0].message.content or "").strip()


async def transcribe_voice(update: Update, context: ContextTypes.DEFAULT_TYPE) -> str:
    message = require_message(update)
    voice = message.voice
    if voice is None:
        raise RuntimeError("Ovozli xabar topilmadi.")
    telegram_file = await context.bot.get_file(voice.file_id)
    with tempfile.TemporaryDirectory(dir=str(make_temp_dir("voice-parent-"))) as tmp_dir:
        audio_path = Path(tmp_dir) / "voice.ogg"
        await telegram_file.download_to_drive(custom_path=str(audio_path))
        with audio_path.open("rb") as audio_file:
            transcription = cast(Any, await asyncio.to_thread(
                lambda: ai_call_with_retries(
                    lambda: openai_client.audio.transcriptions.create(model=OPENAI_TRANSCRIBE_MODEL, file=audio_file)
                )
            ))
        return transcription.text.strip()


async def generate_image(prompt: str) -> Path:
    if not IMAGE_GENERATION_ENABLED:
        raise RuntimeError("Rasm generatsiyasi hozir o'chirilgan.")
    if not OPENAI_IMAGE_API_KEY:
        raise RuntimeError("Rasm generatsiyasi uchun OPENAI_IMAGE_API_KEY kerak.")
    if contains_adult_content(prompt):
        raise ValueError("18+ yoki pornografik rasm so'rovlari qo'llab-quvvatlanmaydi.")
    response = await asyncio.to_thread(
        lambda: image_client.images.generate(
            model=IMAGE_MODEL,
            prompt=(
                "High quality, clean, safe-for-work image. No nudity, no sexual content, no gore. "
                f"User prompt: {prompt}"
            ),
            size=IMAGE_SIZE,
            n=1,
        )
    )
    if not response.data:
        raise RuntimeError("AI rasm qaytarmadi.")
    image_data = response.data[0]
    output_path = make_temp_dir("ai-image-") / "image.png"
    if getattr(image_data, "b64_json", None):
        output_path.write_bytes(base64.b64decode(str(image_data.b64_json)))
        return output_path
    if getattr(image_data, "url", None):
        image_url = str(image_data.url)
        async with httpx.AsyncClient(timeout=60) as client:
            resp = await client.get(image_url)
            resp.raise_for_status()
            output_path.write_bytes(resp.content)
            return output_path
    raise RuntimeError("AI rasm qaytarmadi.")


def image_limit_text(user_id: int) -> str:
    used = get_analytics().user_image_count(user_id)
    return f"Rasm generatsiyasi hozir ochiq. Ishlatilgan: {used}."


def download_media(url: str, media_type: str) -> tuple[Path, str, int | None]:
    try:
        import yt_dlp
        from yt_dlp.utils import DownloadError
    except ImportError as exc:
        raise RuntimeError("yt-dlp o'rnatilmagan. requirements.txt yangilang va deployni qayta qiling.") from exc

    tmp_dir = make_temp_dir("media-")
    output_template = str(tmp_dir / "%(title).80s-%(id)s.%(ext)s")
    has_ffmpeg = shutil.which("ffmpeg") is not None
    options: dict[str, Any] = {
        "outtmpl": output_template,
        "noplaylist": True,
        "quiet": True,
        "noprogress": True,
        "no_warnings": True,
        "max_filesize": MEDIA_MAX_MB * 1024 * 1024,
        "windowsfilenames": True,
        "socket_timeout": 40,
        "retries": 3,
        "fragment_retries": 3,
        "extractor_retries": 3,
    }
    if media_type == "audio":
        options["format"] = "bestaudio[ext=m4a]/bestaudio[ext=mp3]/bestaudio/best[acodec!=none]"
        if has_ffmpeg:
            options["format"] = "bestaudio/best[acodec!=none]/best"
            options["postprocessors"] = [
                {
                    "key": "FFmpegExtractAudio",
                    "preferredcodec": "mp3",
                    "preferredquality": "192",
                }
            ]
    else:
        # Prefer one already-muxed file. Splitting bestvideo+bestaudio needs ffmpeg.
        options["format"] = (
            f"best[ext=mp4][vcodec!=none][acodec!=none][filesize<{MEDIA_MAX_MB}M]/"
            f"best[ext=mp4][vcodec!=none][acodec!=none][filesize_approx<{MEDIA_MAX_MB}M]/"
            f"best[vcodec!=none][acodec!=none][filesize<{MEDIA_MAX_MB}M]/"
            f"best[vcodec!=none][acodec!=none][filesize_approx<{MEDIA_MAX_MB}M]/"
            "best[ext=mp4]/best"
        )
        if has_ffmpeg:
            options["format"] += f"/bestvideo[filesize<{MEDIA_MAX_MB}M]+bestaudio/best"
            options["merge_output_format"] = "mp4"

    with yt_dlp.YoutubeDL(cast(Any, options)) as ydl:
        try:
            info = ydl.extract_info(url, download=True)
        except DownloadError as exc:
            if "requested format is not available" not in str(exc).lower():
                raise
            fallback_options = {**options, "format": "bestaudio/best" if media_type == "audio" else "best[ext=mp4]/best"}
            fallback_options.pop("postprocessors", None)
            fallback_options.pop("merge_output_format", None)
            with yt_dlp.YoutubeDL(cast(Any, fallback_options)) as fallback_ydl:
                info = fallback_ydl.extract_info(url, download=True)
        if info is None:
            raise RuntimeError("Media ma'lumoti qaytmadi.")
        title = str(info.get("title") or "media")
        duration_value = info.get("duration")
        duration = int(duration_value) if duration_value is not None else None
        age_limit = int(info.get("age_limit") or 0)
        if age_limit >= 18 or contains_adult_content(title):
            raise ValueError("18+ media tarqatilmaydi.")
        downloaded = sorted(tmp_dir.glob("*"), key=lambda path: path.stat().st_mtime, reverse=True)
        if not downloaded:
            raise RuntimeError("Yuklangan fayl topilmadi.")
        file_path = downloaded[0]

    if file_path.stat().st_size > MEDIA_MAX_MB * 1024 * 1024:
        raise RuntimeError(f"Fayl {MEDIA_MAX_MB} MB limitdan katta. Qisqaroq video yuboring.")
    return file_path, title, duration


def media_keyboard() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(
        [
            [
                InlineKeyboardButton(
                    "Audio",
                    callback_data="media|audio",
                    icon_custom_emoji_id=AUDIO_BUTTON_CUSTOM_EMOJI_ID,
                    style="success",
                ),
                InlineKeyboardButton(
                    "Video",
                    callback_data="media|video",
                    icon_custom_emoji_id=VIDEO_BUTTON_CUSTOM_EMOJI_ID,
                    style="danger",
                ),
            ]
        ]
    )


def media_done_caption(media_type: str) -> tuple[str, list[MessageEntity]]:
    label = "Audio" if media_type == "audio" else "Video"
    caption = f"{label} @prophoneuzbot dan yuklab olindi {MEDIA_DONE_EMOJI}"
    emoji_offset = len(caption.encode("utf-16-le")) // 2 - len(MEDIA_DONE_EMOJI.encode("utf-16-le")) // 2
    return caption, [
        MessageEntity(
            type=MessageEntity.CUSTOM_EMOJI,
            offset=emoji_offset,
            length=len(MEDIA_DONE_EMOJI.encode("utf-16-le")) // 2,
            custom_emoji_id=MEDIA_DONE_CUSTOM_EMOJI_ID,
        )
    ]


def custom_emoji_html(emoji: str, emoji_id: str) -> str:
    return f'<tg-emoji emoji-id="{emoji_id}">{emoji}</tg-emoji>'


async def reply_html_or_plain(message: Any, text: str) -> None:
    try:
        await message.reply_text(text[:4096], parse_mode=ParseMode.HTML)
    except BadRequest:
        await message.reply_text(text[:4096])


def plain_text(value: str) -> str:
    value = HTML_TAG_RE.sub("", value)
    return html.unescape(value).strip()


def requested_document_format(text: str) -> str:
    lowered = text.lower()
    for extension in ("pptx", "docx", "html", "md", "txt"):
        if extension in lowered:
            return extension
    if "powerpoint" in lowered or "prezentatsiya" in lowered or "презентация" in lowered or "slayd" in lowered:
        return "pptx"
    if "word" in lowered or "referat" in lowered or "konspekt" in lowered:
        return "docx"
    return "pptx"


def safe_filename(title: str, extension: str) -> str:
    cleaned = re.sub(r"[^A-Za-z0-9А-Яа-яЁёЎўҚқҒғҲҳІіЇїЄє _.-]+", "", title).strip()
    cleaned = re.sub(r"\s+", "_", cleaned)[:60] or "ai_material"
    return f"{cleaned}.{extension}"


def parse_slide_outline(text: str) -> tuple[str, list[tuple[str, list[str]]]]:
    lines = [plain_text(line).strip(" -•\t") for line in text.splitlines()]
    lines = [line for line in lines if line]
    title = lines[0].replace("TITLE:", "").replace("MAVZU:", "").strip() if lines else "AI prezentatsiya"
    slides: list[tuple[str, list[str]]] = []
    current_title = ""
    current_items: list[str] = []
    for line in lines[1:]:
        normalized = line.lower()
        if normalized.startswith(("slide:", "slayd:", "слайд:")) or re.match(r"^\d+[.)]\s+", line):
            if current_title:
                slides.append((current_title, current_items[:5]))
            current_title = re.sub(r"^(slide|slayd|слайд):\s*", "", line, flags=re.IGNORECASE)
            current_title = re.sub(r"^\d+[.)]\s*", "", current_title).strip()
            current_items = []
        elif current_title:
            current_items.append(line)
    if current_title:
        slides.append((current_title, current_items[:5]))
    if not slides:
        chunks = lines[1:] or [title]
        for index in range(0, min(len(chunks), 30), 5):
            slides.append((f"Bo'lim {index // 5 + 1}", chunks[index : index + 5]))
    return title, slides[:10]


async def build_learning_material(topic: str) -> str:
    prompt = (
        "Quyidagi mavzu uchun maktab o'qituvchisi va o'quvchilariga mos, tayyor faylga aylantiriladigan "
        "prezentatsiya/o'quv materiali tuz. Faqat oddiy matn qaytar: birinchi qatorda TITLE: mavzu, keyin "
        "8 tagacha SLIDE: sarlavha va har bir slayd ostida 3-5 ta qisqa punkt. Oxirida 5 ta test savoli va "
        f"uyga vazifa qo'sh. Mavzu: {topic}"
    )
    return plain_text(await ask_ai(prompt))


def write_text_material(path: Path, title: str, slides: list[tuple[str, list[str]]], extension: str) -> None:
    if extension == "html":
        body = [f"<h1>{html.escape(title)}</h1>"]
        for slide_title, items in slides:
            body.append(f"<h2>{html.escape(slide_title)}</h2><ul>")
            body.extend(f"<li>{html.escape(item)}</li>" for item in items)
            body.append("</ul>")
        path.write_text("<!doctype html><meta charset='utf-8'>" + "\n".join(body), encoding="utf-8")
        return
    prefix = "#" if extension == "md" else ""
    lines = [f"{prefix} {title}".strip(), ""]
    for slide_title, items in slides:
        lines.extend([f"{prefix * 2} {slide_title}".strip(), ""])
        lines.extend(f"- {item}" for item in items)
        lines.append("")
    path.write_text("\n".join(lines), encoding="utf-8")


def write_docx_material(path: Path, title: str, slides: list[tuple[str, list[str]]]) -> None:
    from docx import Document

    document = Document()
    document.add_heading(title, level=1)
    for slide_title, items in slides:
        document.add_heading(slide_title, level=2)
        for item in items:
            document.add_paragraph(item, style="List Bullet")
    document.save(str(path))


def write_pptx_material(path: Path, title: str, slides: list[tuple[str, list[str]]]) -> None:
    from pptx import Presentation

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_shape = cast(Any, title_slide.shapes.title)
    title_shape.text = title
    cast(Any, title_slide.placeholders[1]).text = "AI yordamida tayyorlandi"
    for slide_title, items in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        cast(Any, slide.shapes.title).text = slide_title
        body = cast(Any, slide.shapes.placeholders[1]).text_frame
        body.clear()
        for index, item in enumerate(items or ["Asosiy fikrlarni o'qituvchi izohi bilan to'ldiring."]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
    presentation.save(str(path))


def create_learning_file(topic: str, material: str, extension: str) -> Path:
    title, slides = parse_slide_outline(material)
    tmp_dir = make_temp_dir("ai-material-")
    path = tmp_dir / safe_filename(title or topic, extension)
    try:
        if extension == "pptx":
            write_pptx_material(path, title, slides)
        elif extension == "docx":
            write_docx_material(path, title, slides)
        else:
            write_text_material(path, title, slides, extension)
    except Exception:
        path = tmp_dir / safe_filename(title or topic, "md")
        write_text_material(path, title, slides, "md")
    return path


async def send_learning_material(update: Update, topic: str) -> None:
    message = require_message(update)
    extension = requested_document_format(topic)
    await message.chat.send_action(ChatAction.UPLOAD_DOCUMENT)
    try:
        material = await build_learning_material(topic)
        file_path = create_learning_file(topic, material, extension)
        with file_path.open("rb") as document_file:
            await message.reply_document(
                document=document_file,
                filename=file_path.name,
                caption="<b>O'quv materiali tayyor</b>",
                parse_mode=ParseMode.HTML,
                read_timeout=60,
                write_timeout=60,
                connect_timeout=20,
                pool_timeout=20,
            )
        record_usage(update, "presentation", text_preview=topic)
    except Exception as exc:
        logger.exception("Learning material generation failed")
        record_usage(update, "presentation", status="error", text_preview=topic, error=str(exc))
        await message.reply_text(
            f"<b>Fayl tayyorlanmadi</b>\n<code>{html.escape(friendly_error(exc))}</code>",
            parse_mode=ParseMode.HTML,
        )


async def start(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message_obj = require_message(update)
    user = require_effective_user(update)
    store = get_analytics()
    args = context_args(context)
    if args and args[0].strip().lower().startswith("love_"):
        code = args[0].strip()[5:].upper()
        session = store.join_love_session(code, user.id, user.full_name)
        if session is None:
            await message_obj.reply_text("<b>Secret chat</b>\n<i>Link noto'g'ri yoki yopilgan.</i>", parse_mode=ParseMode.HTML)
            record_usage(update, "love", status="bad_link")
            return
        if session["partner_id"] and int(session["partner_id"]) != user.id:
            await message_obj.reply_text("<b>Secret chat</b>\n<i>Bu invite allaqachon band.</i>", parse_mode=ParseMode.HTML)
            record_usage(update, "love", status="occupied")
            return
        await message_obj.reply_text(
            "<b>Secret chat ochildi.</b>\n"
            "<i>Bu yerdagi matnlaringiz anonim chat sherigiga yuboriladi. Yopish uchun /love close yozing.</i>",
            parse_mode=ParseMode.HTML,
        )
        try:
            await context.bot.send_message(
                chat_id=int(session["owner_id"]),
                text=(
                    "<b>Secret chatga sherik kirdi.</b>\n"
                    f"Ism: <code>{html.escape(user.full_name)}</code>\n"
                    f"ID: <code>{user.id}</code>\n"
                    "<i>Endi oddiy matn yozsangiz sherikka anonim boradi.</i>"
                ),
                parse_mode=ParseMode.HTML,
            )
        except Exception:
            logger.debug("Love chat owner notice yuborilmadi", exc_info=True)
        record_usage(update, "love", status="joined")
        return
    is_new_user = not store.user_exists(user.id)
    record_usage(update, "start")
    if is_new_user and OWNER_TELEGRAM_ID and user.id != OWNER_TELEGRAM_ID:
        counts = store.subscriber_count()
        username = f"@{user.username}" if user.username else "username yo'q"
        await context.bot.send_message(
            chat_id=OWNER_TELEGRAM_ID,
            text=(
                "Yangi foydalanuvchi qo'shildi\n\n"
                f"Ism: {user.full_name}\n"
                f"Username: {username}\n"
                f"Telegram ID: {user.id}\n"
                f"Jami foydalanuvchi: {counts['total']}"
            )[:4096],
        )
    greeting = f"Salom! Men AI Telegram botman. {START_EMOJI}"
    message = (
        f"{greeting}\n\n"
        "Uzbekcha, ruscha va inglizcha savollarga javob beraman, ovozli xabarni tushunaman, rasm chizaman, "
        "prezentatsiya va o'quv materiallari tayyorlayman, Instagram/YouTube linklaridan audio yoki video tanlashga yordam beraman.\n\n"
        f"Sizning Telegram ID: {user.id}\n"
        f"{image_limit_text(user.id)}\n\n"
        "Privatda savolni oddiy yozing. Guruhlarda meni /ai savol orqali chaqiring."
    )
    entities = [
        MessageEntity(
            type=MessageEntity.BOLD,
            offset=0,
            length=len("Salom! Men AI Telegram botman.".encode("utf-16-le")) // 2,
        ),
        MessageEntity(
            type=MessageEntity.CUSTOM_EMOJI,
            offset=len(greeting.encode("utf-16-le")) // 2 - len(START_EMOJI.encode("utf-16-le")) // 2,
            length=len(START_EMOJI.encode("utf-16-le")) // 2,
            custom_emoji_id=START_CUSTOM_EMOJI_ID,
        ),
    ]
    for attempt in range(3):
        try:
            await message_obj.reply_text(
                message,
                entities=entities,
                read_timeout=60,
                write_timeout=60,
                connect_timeout=30,
                pool_timeout=30,
            )
            return
        except BadRequest:
            await message_obj.reply_text(
                message,
                read_timeout=60,
                write_timeout=60,
                connect_timeout=30,
                pool_timeout=30,
            )
            return
        except (TimedOut, NetworkError):
            if attempt == 2:
                raise
            await asyncio.sleep(2)


async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    record_usage(update, "help")
    await message.reply_text(
        "<b>Yordam</b>\n\n"
        "<b>/ai</b> <i>savol</i> - AI javob\n"
        "<b>/image</b> <i>prompt</i> - rasm generatsiyasi\n"
        "<b>/present</b> <i>mavzu</i> - prezentatsiya yoki o'quv fayli\n"
        "<b>/media audio</b> <i>LINK</i> - YouTube/Instagram audio\n"
        "<b>/media video</b> <i>LINK</i> - YouTube/Instagram video\n"
        "<b>/radar</b> - trendlar\n\n"
        "<i>Privatda matn yoki ovoz yuborsangiz ham AI javob beradi. Guruhlarda faqat /ai orqali ishlayman.</i>",
        parse_mode=ParseMode.HTML,
    )


async def radar_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    user = require_effective_user(update)
    if not is_owner(user.id):
        await message.reply_text("<b>Radar</b>\n<i>Bu bo'lim faqat bot egasi uchun.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "radar", status="denied")
        return
    counts = get_analytics().subscriber_count()
    top_queries = get_analytics().top_queries(days=7, limit=10)
    lines = [
        "<b>Bot radari</b>",
        f"Jami foydalanuvchi: {counts['total']}",
        f"Bugun qo'shilgan: {counts['new_today']}",
        f"7 kunda qo'shilgan: {counts['new_week']}",
        f"Bugungi faol: {counts['active_today']}",
        f"7 kunlik faol: {counts['active_week']}",
        "",
        "<b>Ko'p qidirilgan/so'ralganlar</b>",
    ]
    lines.extend(f"- {html.escape(item)}" for item in top_queries) if top_queries else lines.append("<i>Hali yetarli so'rov yo'q</i>")
    record_usage(update, "radar")
    await message.reply_text("\n".join(lines)[:4096], parse_mode=ParseMode.HTML)


async def report_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    user = require_effective_user(update)
    if not is_owner(user.id):
        await message.reply_text("<b>Hisobot</b>\n<i>Bu bo'lim faqat bot egasi uchun.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "report", status="denied")
        return
    report = get_analytics().build_report("Oxirgi 7 kunlik Telegram AI bot hisoboti", start_at=datetime.now() - timedelta(days=7))
    record_usage(update, "report")
    await message.reply_text(f"<pre>{html.escape(report[:4000])}</pre>", parse_mode=ParseMode.HTML)


async def love_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    user = require_effective_user(update)
    args = context_args(context)
    store = get_analytics()
    if args and args[0].strip().lower() in {"close", "stop", "yop"}:
        session = store.close_love_session(user.id)
        if session is None:
            await message.reply_text("<b>Secret chat</b>\n<i>Aktiv chat topilmadi.</i>", parse_mode=ParseMode.HTML)
            return
        other_id = int(session["partner_id"] or 0) if user.id == int(session["owner_id"]) else int(session["owner_id"])
        await message.reply_text("<b>Secret chat yopildi.</b>", parse_mode=ParseMode.HTML)
        if other_id:
            try:
                await context.bot.send_message(chat_id=other_id, text="<b>Secret chat yopildi.</b>", parse_mode=ParseMode.HTML)
            except Exception:
                logger.debug("Love chat close notice yuborilmadi", exc_info=True)
        record_usage(update, "love", status="closed")
        return

    if user.id == LOVE_OWNER_ID and (not args or args[0].strip().lower() in {"new", "invite", "link"}):
        code = make_love_code()
        store.create_love_invite(user.id, code)
        bot_user = await context.bot.get_me()
        link = f"https://t.me/{bot_user.username}?start=love_{code}" if bot_user.username else ""
        lines = [
            "<b>Secret chat invite tayyor</b>",
            f"Kod: <code>{code}</code>",
            "",
            "Sherikka shuni yuboring:",
            f"<code>/love {code}</code>",
        ]
        if link:
            lines.extend(["", "Yoki link:", f"<code>{html.escape(link)}</code>"])
        lines.append("")
        lines.append("<i>Chatga kirgan odamga xabarlari anonim sherikka yuborilishi aytiladi.</i>")
        await message.reply_text("\n".join(lines), parse_mode=ParseMode.HTML, disable_web_page_preview=True)
        record_usage(update, "love", status="invite_created")
        return

    if not args:
        if user.id == LOVE_OWNER_ID:
            await message.reply_text(
                "<b>Secret chat</b>\n<code>/love new</code> - yangi invite\n<code>/love close</code> - chatni yopish",
                parse_mode=ParseMode.HTML,
            )
        return

    code = args[0].strip().upper()
    session = store.join_love_session(code, user.id, user.full_name)
    if session is None:
        await message.reply_text("<b>Secret chat</b>\n<i>Kod noto'g'ri yoki yopilgan.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "love", status="bad_code")
        return
    if session["partner_id"] and int(session["partner_id"]) != user.id:
        await message.reply_text("<b>Secret chat</b>\n<i>Bu invite allaqachon band.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "love", status="occupied")
        return
    await message.reply_text(
        "<b>Secret chat ochildi.</b>\n"
        "<i>Bu yerdagi matnlaringiz anonim chat sherigiga yuboriladi. Yopish uchun /love close yozing.</i>",
        parse_mode=ParseMode.HTML,
    )
    try:
        await context.bot.send_message(
            chat_id=int(session["owner_id"]),
            text=(
                "<b>Secret chatga sherik kirdi.</b>\n"
                f"Ism: <code>{html.escape(user.full_name)}</code>\n"
                f"ID: <code>{user.id}</code>\n"
                "<i>Endi oddiy matn yozsangiz sherikka anonim boradi.</i>"
            ),
            parse_mode=ParseMode.HTML,
        )
    except Exception:
        logger.debug("Love chat owner notice yuborilmadi", exc_info=True)
    record_usage(update, "love", status="joined")


async def maybe_handle_love_message(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str) -> bool:
    if is_group_chat(update):
        return False
    message = require_message(update)
    user = require_effective_user(update)
    session = get_analytics().get_love_session_by_user(user.id)
    if session is None or not session["partner_id"]:
        return False
    owner_id = int(session["owner_id"])
    partner_id = int(session["partner_id"])
    if user.id == owner_id:
        target_id = partner_id
        prefix = "<b>Anonim xabar</b>"
    elif user.id == partner_id:
        target_id = owner_id
        partner_name = html.escape(str(session["partner_name"] or "Sherik"))
        prefix = f"<b>Secret chat</b>\n<i>{partner_name}:</i>"
    else:
        return False
    try:
        await context.bot.send_message(
            chat_id=target_id,
            text=f"{prefix}\n{html.escape(text)[:3800]}",
            parse_mode=ParseMode.HTML,
        )
        await message.reply_text("<i>Yuborildi.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "love", text_preview=text)
    except Exception as exc:
        logger.exception("Love chat relay failed")
        await message.reply_text(
            f"<b>Secret chat xabari yuborilmadi</b>\n<code>{html.escape(friendly_error(exc))}</code>",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "love", status="error", text_preview=text, error=str(exc))
    return True


async def ai_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    text = " ".join(context_args(context)).strip()
    if not text:
        await message.reply_text(
            "<b>Savolni ham yozing</b>\n<i>Masalan:</i> <code>/ai nima yordam bera olasan?</code>",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "ai", status="empty")
        return
    await answer_text(update, text, action="ai")


async def image_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    prompt = " ".join(context_args(context)).strip()
    if not prompt:
        await message.reply_text(
            "<b>Rasm uchun prompt yozing</b>\n<i>Masalan:</i> <code>/image futuristik Toshkent</code>",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "image", status="empty")
        return
    await handle_image_request(update, prompt)


async def present_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    topic = " ".join(context_args(context)).strip()
    if not topic:
        await message.reply_text(
            "<b>Mavzu yozing</b>\n<i>Masalan:</i> <code>/present Amir Temur haqida 8 slayd pptx</code>",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "presentation", status="empty")
        return
    await send_learning_material(update, topic)


async def media_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    args = context_args(context)
    if not MEDIA_DOWNLOAD_ENABLED:
        await message.reply_text("<b>Media yuklash</b>\n<i>Hozir o'chirilgan.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "media", status="disabled")
        return
    if len(args) < 2:
        await message.reply_text(
            "<b>Format</b>\n<code>/media audio LINK</code>\nyoki\n<code>/media video LINK</code>",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "media", status="empty")
        return
    media_type = args[0].strip().lower()
    url = args[1].strip()
    if media_type not in {"audio", "video"}:
        await message.reply_text(
            "<b>Format xato</b>\nBirinchi so'z <code>audio</code> yoki <code>video</code> bo'lishi kerak.",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "media", status="bad_type", text_preview=" ".join(args))
        return
    await handle_media_download(update, url, media_type)


async def kinotop_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    user = update.effective_user
    if not user or user.id != KINOTOP_OWNER_ID:
        return
    args = context_args(context)
    if not args:
        await message.reply_text("<b>Link yuboring</b>\n<code>/kinotop https://site.com/video</code>", parse_mode=ParseMode.HTML)
        record_usage(update, "kinotop", status="empty")
        return
    if not MEDIA_DOWNLOAD_ENABLED:
        await message.reply_text("<b>Media yuklash</b>\n<i>Hozir o'chirilgan.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "kinotop", status="disabled")
        return
    if not KINOTOP_ALLOWED_DOMAINS:
        await message.reply_text(
            "<b>Kinotop sozlanmagan</b>\n<code>KINOTOP_ALLOWED_DOMAINS</code> ichiga domenlarni vergul bilan yozing.",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "kinotop", status="not_configured")
        return
    url = args[0].strip()
    if not URL_RE.fullmatch(url):
        await message.reply_text("<b>Link noto'g'ri</b>\nTo'liq <code>https://...</code> link yuboring.", parse_mode=ParseMode.HTML)
        record_usage(update, "kinotop", status="bad_url", text_preview=url)
        return
    if not is_kinotop_allowed_url(url):
        await message.reply_text("<b>Domen ruxsatda yo'q</b>\nBu link whitelist ichidagi saytdan emas.", parse_mode=ParseMode.HTML)
        record_usage(update, "kinotop", status="domain_blocked", text_preview=get_url_domain(url))
        return
    await handle_media_download(update, url, "video", allow_custom_url=True)


async def admin_status_command(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    user = update.effective_user
    if not user or user.id != KINOTOP_OWNER_ID:
        return
    lines = [
        "<b>Admin status</b>",
        f"AI provider: <code>{html.escape(AI_PROVIDER)}</code>",
        f"AI model: <code>{html.escape(OPENAI_TEXT_MODEL)}</code>",
        f"AI key: <code>{'bor' if bool(AI_API_KEY) else 'yoq'}</code>",
        f"AI base URL: <code>{html.escape(AI_BASE_URL or 'default')}</code>",
        f"Media download: <code>{'yoqilgan' if MEDIA_DOWNLOAD_ENABLED else 'ochirilgan'}</code>",
        f"Media max MB: <code>{MEDIA_MAX_MB}</code>",
        f"Kinotop owner: <code>{KINOTOP_OWNER_ID}</code>",
        f"Kinotop domenlar: <code>{len(KINOTOP_ALLOWED_DOMAINS)}</code>",
    ]
    if KINOTOP_ALLOWED_DOMAINS:
        lines.append(f"Domenlar: <code>{html.escape(', '.join(sorted(KINOTOP_ALLOWED_DOMAINS)))}</code>")
    await message.reply_text("\n".join(lines), parse_mode=ParseMode.HTML)


async def media_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    query = cast(Any, update.callback_query)
    if query is None:
        return
    await query.answer()
    _, media_type = query.data.split("|", 1)
    user_data = cast(dict[str, Any], context.user_data or {})
    url = user_data.get("pending_media_url", "")
    if not url:
        await query.message.reply_text("<b>Link topilmadi</b>\n<i>Iltimos, linkni qayta yuboring.</i>", parse_mode=ParseMode.HTML)
        return
    await handle_media_download(update, url, media_type, from_callback=True)


async def handle_media_download(
    update: Update,
    url: str,
    media_type: str,
    from_callback: bool = False,
    allow_custom_url: bool = False,
) -> None:
    message = require_message(update)
    if not allow_custom_url and not MEDIA_URL_RE.match(url):
        await message.reply_text("<b>Link noto'g'ri</b>\nFaqat <i>YouTube</i> yoki <i>Instagram</i> link yuboring.", parse_mode=ParseMode.HTML)
        record_usage(update, "media", status="bad_url", text_preview=url)
        return
    if contains_adult_content(url):
        await message.reply_text("<b>Cheklov</b>\n<i>18+ materiallar tarqatilmaydi.</i>", parse_mode=ParseMode.HTML)
        record_usage(update, "media", status="adult_blocked", text_preview=url)
        return
    await message.chat.send_action(ChatAction.UPLOAD_DOCUMENT)
    try:
        file_path, title, duration = await asyncio.to_thread(download_media, url, media_type)
        with file_path.open("rb") as media_file:
            caption, caption_entities = media_done_caption(media_type)
            if media_type == "audio":
                await message.reply_audio(
                    audio=media_file,
                    title=title[:64],
                    filename=file_path.name,
                    duration=duration,
                    caption=caption,
                    caption_entities=caption_entities,
                    read_timeout=60,
                    write_timeout=60,
                    connect_timeout=20,
                    pool_timeout=20,
                )
            else:
                await message.reply_document(
                    document=media_file,
                    filename=file_path.name,
                    caption=caption,
                    caption_entities=caption_entities,
                    read_timeout=60,
                    write_timeout=60,
                    connect_timeout=20,
                    pool_timeout=20,
                )
    except Exception as exc:
        logger.exception("Media download failed")
        record_usage(update, "media", status="error", text_preview=url, error=str(exc))
        await message.reply_text(
            f"<b>Media tayyorlanmadi</b>\n<code>{html.escape(friendly_error(exc))}</code>",
            parse_mode=ParseMode.HTML,
        )
        return
    record_usage(update, "media", text_preview=url)
    if from_callback:
        try:
            query = cast(Any, update.callback_query)
            if query is not None:
                await query.edit_message_reply_markup(reply_markup=None)
        except Exception:
            logger.debug("Callback markupni o'chirib bo'lmadi", exc_info=True)


async def answer_text(update: Update, text: str, action: str = "text") -> None:
    message = require_message(update)
    if contains_adult_content(text):
        await message.reply_text(
            "<b>Bu mavzuda yordam bera olmayman</b>\n<i>Xavfsiz, ta'limiy yoki ijodiy mavzu tanlang.</i>",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, action, status="adult_blocked", text_preview=text)
        return
    await message.chat.send_action(ChatAction.TYPING)
    try:
        answer = await ask_ai(text)
    except Exception as exc:
        logger.exception("AI request failed")
        record_usage(update, action, status="error", text_preview=text, error=str(exc))
        await message.reply_text(
            f"<b>AI javob bera olmadi</b>\n<code>{html.escape(friendly_error(exc))}</code>",
            parse_mode=ParseMode.HTML,
        )
        return
    record_usage(update, action, text_preview=text, response_chars=len(answer))
    await reply_html_or_plain(message, answer)


async def handle_image_request(update: Update, prompt: str) -> None:
    message = require_message(update)
    user_id = require_effective_user(update).id
    await message.chat.send_action(ChatAction.UPLOAD_PHOTO)
    try:
        image_path = await generate_image(prompt)
        with image_path.open("rb") as image_file:
            await message.reply_photo(
                photo=image_file,
                caption=f"<b>Rasm tayyor</b>\n{html.escape(image_limit_text(user_id))}"[:1000],
                parse_mode=ParseMode.HTML,
            )
    except Exception as exc:
        logger.exception("Image generation failed")
        record_usage(update, "image", status="error", text_preview=prompt, error=str(exc))
        await message.reply_text(
            f"<b>Rasm tayyorlanmadi</b>\n<code>{html.escape(friendly_error(exc))}</code>",
            parse_mode=ParseMode.HTML,
        )
        return
    record_usage(update, "image", text_preview=prompt)


async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    if is_group_chat(update):
        await message.reply_text(
            "<b>Ovozli xabar</b>\nGuruhlarda AI uchun <code>/ai savol</code> yozing. "
            "<i>Privatda voice yuborsangiz javob beraman.</i>",
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "voice", status="group_ignored")
        return
    await message.chat.send_action(ChatAction.TYPING)
    try:
        text = await transcribe_voice(update, context)
    except Exception as exc:
        logger.exception("Voice transcription failed")
        record_usage(update, "voice", status="error", error=str(exc))
        await message.reply_text(
            f"<b>Ovozni matnga aylantirib bo'lmadi</b>\n<code>{html.escape(friendly_error(exc))}</code>",
            parse_mode=ParseMode.HTML,
        )
        return
    record_usage(update, "voice", text_preview=text)
    await answer_text(update, text, action="ai")


async def handle_text(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    message = require_message(update)
    text = (message.text or "").strip()
    if await maybe_handle_love_message(update, context, text):
        return
    urls = URL_RE.findall(text)
    if urls and MEDIA_URL_RE.match(urls[0]):
        user_data = cast(dict[str, Any], context.user_data or {})
        user_data["pending_media_url"] = urls[0]
        await message.reply_text(
            "<b>Nimani olishni tanlang?</b>",
            reply_markup=media_keyboard(),
            parse_mode=ParseMode.HTML,
        )
        record_usage(update, "media", status="choice", text_preview=urls[0])
        return
    if is_group_chat(update):
        return
    if looks_like_presentation_request(text):
        await send_learning_material(update, text)
        return
    if looks_like_image_request(text):
        await handle_image_request(update, text)
        return
    await answer_text(update, text, action="text")


async def greet_new_members(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if not update.message or not update.message.new_chat_members:
        return
    bot_id = context.bot.id
    if any(member.id == bot_id for member in update.message.new_chat_members):
        record_usage(update, "group_join")
        await update.message.reply_text(
            "<b>Salom, men AI yordamchi botman.</b>\n\n"
            "Guruhda tartibli ishlash uchun meni <code>/ai savol</code> orqali chaqiring.\n"
            "<i>Uzbekcha, ruscha va inglizcha javob beraman; 18+ materiallarni tarqatmayman.</i>",
            parse_mode=ParseMode.HTML,
        )


async def error_handler(update: object, context: ContextTypes.DEFAULT_TYPE) -> None:
    if isinstance(context.error, (TimedOut, NetworkError)):
        logger.warning("Telegram network timeout/error: %s", context.error)
        return

    logger.exception("Telegram update failed", exc_info=context.error)
    if isinstance(update, Update) and update.effective_message:
        try:
            await update.effective_message.reply_text(
                "<b>Kechirasiz, bot ichida xatolik chiqdi.</b>\n<i>Birozdan keyin yana urinib ko'ring.</i>",
                parse_mode=ParseMode.HTML,
            )
        except Exception:
            logger.exception("Error xabarini yuborib bo'lmadi")


def build_application() -> Application:
    application = (
        Application.builder()
        .token(TELEGRAM_BOT_TOKEN)
        .connect_timeout(30)
        .read_timeout(60)
        .write_timeout(60)
        .pool_timeout(30)
        .get_updates_connect_timeout(30)
        .get_updates_read_timeout(60)
        .get_updates_write_timeout(60)
        .get_updates_pool_timeout(30)
        .media_write_timeout(120)
        .post_init(start_background_tasks)
        .post_shutdown(stop_background_tasks)
        .build()
    )
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CommandHandler("radar", radar_command))
    application.add_handler(CommandHandler("necha_yulduz", radar_command))
    application.add_handler(CommandHandler("report", report_command))
    application.add_handler(CommandHandler("love", love_command))
    application.add_handler(CommandHandler("ai", ai_command))
    application.add_handler(CommandHandler("image", image_command))
    application.add_handler(CommandHandler("rasm", image_command))
    application.add_handler(CommandHandler("present", present_command))
    application.add_handler(CommandHandler("prezentatsiya", present_command))
    application.add_handler(CommandHandler("media", media_command))
    application.add_handler(CommandHandler("yt_ol", media_command))
    application.add_handler(CommandHandler("kinotop", kinotop_command))
    application.add_handler(CommandHandler("admin_status", admin_status_command))
    application.add_handler(CallbackQueryHandler(media_callback, pattern=r"^media\|"))
    application.add_handler(MessageHandler(filters.StatusUpdate.NEW_CHAT_MEMBERS, greet_new_members))
    application.add_handler(MessageHandler(filters.VOICE, handle_voice))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    application.add_error_handler(error_handler)
    return application


app = FastAPI()


def clean_base_url(value: str) -> str:
    value = value.strip().rstrip("/")
    if not value:
        return ""
    if not value.startswith(("http://", "https://")):
        value = f"https://{value}"
    return value


def configured_base_url(request: Request | None = None) -> str:
    for key in ("PUBLIC_BASE_URL", "VERCEL_PROJECT_PRODUCTION_URL", "VERCEL_URL"):
        base_url = clean_base_url(os.getenv(key, ""))
        if base_url:
            return base_url
    if request:
        return clean_base_url(str(request.base_url))
    return ""


def webhook_url_for(request: Request | None = None, explicit_url: str | None = None) -> str:
    if explicit_url:
        return clean_base_url(explicit_url)
    base_url = configured_base_url(request)
    if not base_url:
        raise RuntimeError("PUBLIC_BASE_URL yoki VERCEL_URL topilmadi.")
    return f"{base_url}/telegram-webhook"


async def get_bot_application() -> Application:
    global bot_app, bot_init_error, analytics

    if bot_app:
        return bot_app

    async with bot_lock:
        if bot_app:
            return bot_app
        try:
            require_config()
            analytics = AnalyticsStore(ANALYTICS_DB_FILE)
            bot_app = build_application()
            await bot_app.initialize()
            await setup_bot_commands(bot_app)
            if not bot_app.running:
                await bot_app.start()
            bot_init_error = ""
            return bot_app
        except Exception as exc:
            bot_init_error = str(exc)[:700]
            raise


async def ensure_webhook(request: Request | None = None, explicit_url: str | None = None) -> dict[str, object]:
    global webhook_ready

    async with webhook_lock:
        application = await get_bot_application()
        url = webhook_url_for(request, explicit_url)
        info = await application.bot.get_webhook_info()
        if info.url == url:
            webhook_ready = True
            return {"ok": True, "webhook_url": url, "already_set": True}
        ok = await application.bot.set_webhook(url=url, allowed_updates=Update.ALL_TYPES, drop_pending_updates=False)
        webhook_ready = bool(ok)
        return {"ok": ok, "webhook_url": url, "already_set": False}


async def force_webhook(request: Request | None = None, explicit_url: str | None = None) -> dict[str, object]:
    global webhook_ready

    async with webhook_lock:
        application = await get_bot_application()
        url = webhook_url_for(request, explicit_url)
        await application.bot.delete_webhook(drop_pending_updates=True)
        ok = await application.bot.set_webhook(
            url=url,
            allowed_updates=Update.ALL_TYPES,
            drop_pending_updates=True,
            max_connections=40,
        )
        info = await application.bot.get_webhook_info()
        webhook_ready = bool(ok)
        return {
            "ok": ok,
            "webhook_url": url,
            "current_url": info.url,
            "pending_update_count": info.pending_update_count,
            "last_error_date": info.last_error_date.isoformat() if info.last_error_date else None,
            "last_error_message": info.last_error_message,
        }


@app.get("/")
async def health_check(request: Request) -> dict[str, object]:
    return {
        "status": "ok",
        "service": "telegram-ai-bot",
        "entrypoint": "bot.py",
        "base_url": configured_base_url(request),
        "webhook_url": webhook_url_for(request),
        "next": ["/status", "/ai-health", "/setup-webhook", "/webhook-info"],
    }


@app.get("/status")
async def status() -> dict[str, object]:
    missing_config = []
    if not TELEGRAM_BOT_TOKEN:
        missing_config.append("TELEGRAM_BOT_TOKEN")
    if not AI_API_KEY:
        missing_config.append("GROQ_API_KEY yoki OPENAI_API_KEY")
    return {
        "status": "ok",
        "entrypoint": "bot.py",
        "telegram_token": bool(TELEGRAM_BOT_TOKEN),
        "ai_provider": AI_PROVIDER,
        "ai_model": OPENAI_TEXT_MODEL,
        "transcribe_model": OPENAI_TRANSCRIBE_MODEL,
        "ai_key": bool(AI_API_KEY),
        "image_key": bool(OPENAI_IMAGE_API_KEY),
        "image_generation": IMAGE_GENERATION_ENABLED,
        "media_download": MEDIA_DOWNLOAD_ENABLED,
        "analytics_db": ANALYTICS_DB_FILE,
        "missing_config": missing_config,
        "bot_initialized": bot_app is not None,
        "bot_init_error": bot_init_error,
    }


@app.get("/ai-health")
async def ai_health() -> dict[str, object]:
    try:
        answer = await ask_ai("Faqat bitta so'z bilan javob ber: OK")
        return {"ok": True, "provider": AI_PROVIDER, "model": OPENAI_TEXT_MODEL, "answer": answer[:100]}
    except Exception as exc:
        logger.exception("AI health check failed")
        return {"ok": False, "provider": AI_PROVIDER, "model": OPENAI_TEXT_MODEL, "error": str(exc)}


@app.get("/setup-webhook")
async def setup_webhook(request: Request, url: str | None = None) -> dict[str, object]:
    try:
        return await ensure_webhook(request, url)
    except Exception as exc:
        logger.exception("Webhook setup failed")
        return {
            "ok": False,
            "error": str(exc)[:700],
            "hint": "Vercel Environment Variables ichida TELEGRAM_BOT_TOKEN va AI key nomlarini tekshiring.",
        }


@app.get("/force-webhook")
async def force_webhook_endpoint(request: Request, url: str | None = None) -> dict[str, object]:
    try:
        return await force_webhook(request, url)
    except Exception as exc:
        logger.exception("Force webhook setup failed")
        return {"ok": False, "error": str(exc)[:700]}


@app.get("/webhook-info")
async def webhook_info() -> dict[str, object]:
    try:
        application = await get_bot_application()
        info = await application.bot.get_webhook_info()
        return {
            "ok": True,
            "url": info.url,
            "pending_update_count": info.pending_update_count,
            "last_error_date": info.last_error_date.isoformat() if info.last_error_date else None,
            "last_error_message": info.last_error_message,
        }
    except Exception as exc:
        logger.exception("Webhook info failed")
        return {"ok": False, "error": str(exc)[:700]}


@app.get("/cron/weekly")
async def weekly_cron() -> dict[str, object]:
    get_analytics()
    application = await get_bot_application()
    await maybe_send_reports(application)
    return {"ok": True, "message": "weekly reports checked"}


@app.post("/telegram-webhook")
async def telegram_webhook(request: Request) -> dict[str, bool]:
    try:
        application = await get_bot_application()
        payload = await request.json()
        update = Update.de_json(payload, application.bot)
        if update is None:
            raise HTTPException(status_code=400, detail="Invalid Telegram update")
        await application.process_update(update)
        return {"ok": True}
    except HTTPException:
        raise
    except Exception as exc:
        logger.exception("Webhook update failed")
        raise HTTPException(status_code=500, detail=str(exc)) from exc


if __name__ == "__main__":
    require_config()
    analytics = AnalyticsStore(ANALYTICS_DB_FILE)
    polling_app = build_application()
    logger.info("AI bot ishga tushdi.")
    polling_app.run_polling(allowed_updates=Update.ALL_TYPES)
